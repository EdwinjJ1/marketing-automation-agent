"""
文档解析模块
支持 PDF, Word, PowerPoint, Markdown 等格式的解析
"""

import os
from typing import List, Dict, Any, Optional
from pathlib import Path
import markdown
from bs4 import BeautifulSoup
from docx import Document as DocxDocument


class DocumentParser:
    """通用文档解析器"""

    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.md', '.txt', '.pptx']

    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        解析文档并提取结构化内容

        Args:
            file_path: 文档路径

        Returns:
            包含以下键的字典:
            - text: 提取的文本内容
            - metadata: 文档元数据
            - sections: 文档章节
            - images: 文档中的图片列表
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        suffix = path.suffix.lower()

        if suffix not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {suffix}")

        # 根据文件类型选择解析方法
        if suffix == '.pdf':
            return self._parse_pdf(file_path)
        elif suffix == '.docx':
            return self._parse_docx(file_path)
        elif suffix == '.md':
            return self._parse_markdown(file_path)
        elif suffix == '.txt':
            return self._parse_text(file_path)
        elif suffix == '.pptx':
            return self._parse_pptx(file_path)

    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """解析PDF文档"""
        try:
            import pymupdf4llm
            md_text = pymupdf4llm.to_markdown(file_path)
            return self._parse_markdown_content(md_text, file_path)
        except ImportError:
            # 降级方案：使用 PyMuPDF
            import fitz
            doc = fitz.open(file_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return {
                "text": text,
                "metadata": {"source": file_path, "type": "pdf"},
                "sections": self._extract_sections(text),
                "images": []
            }

    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """解析Word文档"""
        doc = DocxDocument(file_path)

        # 提取段落
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        text = "\n\n".join(paragraphs)

        # 提取图片
        images = []
        for shape in doc.inline_shapes:
            if shape.graphic:
                images.append({
                    "type": "embedded",
                    "size": (shape.width, shape.height)
                })

        # 提取元数据
        metadata = {
            "source": file_path,
            "type": "docx",
            "author": doc.core_properties.author,
            "title": doc.core_properties.title,
            "created": str(doc.core_properties.created)
        }

        return {
            "text": text,
            "metadata": metadata,
            "sections": self._extract_sections(text),
            "images": images
        }

    def _parse_markdown(self, file_path: str) -> Dict[str, Any]:
        """解析Markdown文档"""
        with open(file_path, 'r', encoding='utf-8') as f:
            md_text = f.read()

        return self._parse_markdown_content(md_text, file_path)

    def _parse_markdown_content(self, md_text: str, source: str) -> Dict[str, Any]:
        """解析Markdown内容"""
        # 转换为HTML以便提取结构
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, 'html.parser')

        # 提取标题作为章节
        sections = []
        for heading in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
            level = int(heading.name[1])
            sections.append({
                "level": level,
                "title": heading.get_text().strip(),
                "content": heading.get_text().strip()
            })

        # 提取纯文本
        text = soup.get_text()

        return {
            "text": text,
            "metadata": {"source": source, "type": "markdown"},
            "sections": sections,
            "images": self._extract_markdown_images(md_text)
        }

    def _parse_text(self, file_path: str) -> Dict[str, Any]:
        """解析纯文本文档"""
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()

        return {
            "text": text,
            "metadata": {"source": file_path, "type": "text"},
            "sections": self._extract_sections(text),
            "images": []
        }

    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """解析PowerPoint文档"""
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text = []

        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            slides_text.append("\n".join(slide_text))

        text = "\n\n---\n\n".join(slides_text)

        return {
            "text": text,
            "metadata": {
                "source": file_path,
                "type": "pptx",
                "slides": len(prs.slides)
            },
            "sections": [
                {"level": 1, "title": f"幻灯片 {i+1}", "content": text}
                for i, text in enumerate(slides_text)
            ],
            "images": []
        }

    def _extract_sections(self, text: str) -> List[Dict[str, Any]]:
        """从文本中提取章节"""
        import re

        sections = []
        # 匹配常见的标题格式（# 标题，数字标题等）
        patterns = [
            r'^#{1,6}\s+(.+)$',  # Markdown标题
            r'^[第\s\d+]+章\s+(.+)$',  # 中文章节
            r'^\d+\.\d+\s+(.+)$',  # 数字编号
        ]

        lines = text.split('\n')
        for line in lines:
            for pattern in patterns:
                match = re.match(pattern, line.strip())
                if match:
                    sections.append({
                        "title": match.group(1),
                        "content": line.strip()
                    })
                    break

        return sections

    def _extract_markdown_images(self, md_text: str) -> List[Dict[str, str]]:
        """从Markdown中提取图片链接"""
        import re

        images = []
        pattern = r'!\[([^\]]*)\]\(([^)]+)\)'
        matches = re.findall(pattern, md_text)

        for alt, url in matches:
            images.append({
                "alt": alt,
                "url": url,
                "type": "external" if url.startswith('http') else "local"
            })

        return images

    def extract_key_info(self, parsed_doc: Dict[str, Any]) -> Dict[str, Any]:
        """
        从解析的文档中提取关键信息
        用于营销内容生成
        """
        text = parsed_doc["text"]

        # 使用LLM提取关键信息（简化版）
        key_info = {
            "title": self._extract_title(parsed_doc),
            "summary": self._extract_summary(text),
            "key_features": self._extract_features(text),
            "target_audience": self._guess_audience(text),
            "tone": self._guess_tone(text),
        }

        return key_info

    def _extract_title(self, doc: Dict[str, Any]) -> str:
        """提取文档标题"""
        metadata = doc.get("metadata", {})

        # 优先使用元数据中的标题
        if metadata.get("title"):
            return metadata["title"]

        # 从第一个一级标题提取
        sections = doc.get("sections", [])
        for section in sections:
            if section.get("level") == 1:
                return section["title"]

        # 使用文件名
        return Path(metadata["source"]).stem

    def _extract_summary(self, text: str, max_length: int = 200) -> str:
        """提取文档摘要"""
        # 简单实现：取前N个字符
        # 实际应该使用LLM来生成摘要
        if len(text) <= max_length:
            return text.strip()

        summary = text[:max_length]
        # 找到最后一个完整的句子
        last_period = summary.rfind('。')
        if last_period > max_length * 0.7:
            summary = summary[:last_period + 1]

        return summary.strip()

    def _extract_features(self, text: str) -> List[str]:
        """提取关键特性"""
        # 简化版：查找常见的特性关键词
        # 实际应该使用LLM + NLP
        features = []

        # 查找列表项
        import re
        list_items = re.findall(r'^[\s]*[-*•]\s+(.+)$', text, re.MULTILINE)
        features = list_items[:5]  # 取前5个

        return features

    def _guess_audience(self, text: str) -> str:
        """推测目标受众"""
        # 简化版：根据关键词推测
        # 实际应该使用LLM分析
        text_lower = text.lower()

        if any(word in text_lower for word in ['developer', '程序员', '开发']):
            return "开发者"
        elif any(word in text_lower for word in ['business', '企业', '商业']):
            return "企业用户"
        elif any(word in text_lower for word in ['student', '学生', '学习']):
            return "学生"
        else:
            return "大众"

    def _guess_tone(self, text: str) -> str:
        """推测文档语调"""
        # 简化版
        # 实际应该使用LLM分析
        return "专业"
